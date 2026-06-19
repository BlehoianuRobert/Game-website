from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# Brand colors
COLOR_BG        = RGBColor(0x0D, 0x0D, 0x1A)   # near-black navy
COLOR_ACCENT    = RGBColor(0x00, 0xE5, 0xFF)   # electric cyan
COLOR_ACCENT2   = RGBColor(0x7C, 0x3A, 0xED)   # purple
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT     = RGBColor(0xB0, 0xC4, 0xDE)
COLOR_CARD      = RGBColor(0x1A, 0x1A, 0x2E)   # card bg
COLOR_CARD2     = RGBColor(0x16, 0x21, 0x3E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout

# ─────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────

def add_bg(slide, color=COLOR_BG):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_rect(slide, l, t, w, h, fill=COLOR_CARD, line=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=COLOR_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_accent_bar(slide, t=0.85, h=0.07, color=COLOR_ACCENT):
    bar = slide.shapes.add_shape(1, 0, Inches(t), prs.slide_width, Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def add_bullet_box(slide, items, l, t, w, h, title=None, title_color=COLOR_ACCENT, dot_color=COLOR_ACCENT2):
    add_rect(slide, l, t, w, h, fill=COLOR_CARD, line=COLOR_ACCENT2)
    y = t + 0.18
    if title:
        add_text(slide, title, l+0.15, t+0.08, w-0.3, 0.35, size=13, bold=True, color=title_color)
        y = t + 0.45
    for item in items:
        add_text(slide, f"▸  {item}", l+0.2, y, w-0.4, 0.32, size=11, color=COLOR_LIGHT)
        y += 0.33

def add_tag(slide, text, l, t, color=COLOR_ACCENT2):
    w = len(text)*0.085 + 0.3
    add_rect(slide, l, t, w, 0.28, fill=color)
    add_text(slide, text, l+0.08, t+0.02, w-0.1, 0.26, size=10, bold=True, color=COLOR_WHITE)

# ─────────────────────────────────────────────────────────
# SLIDE 1 — Title / Team
# ─────────────────────────────────────────────────────────
def slide_1():
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)

    # gradient overlay strip left
    g = slide.shapes.add_shape(1, 0, 0, Inches(5), prs.slide_height)
    g.fill.solid()
    g.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x1F)
    g.line.fill.background()

    add_accent_bar(slide, t=0.82, h=0.06, color=COLOR_ACCENT)
    add_accent_bar(slide, t=0.88, h=0.03, color=COLOR_ACCENT2)

    # GRIDFORGE big title
    add_text(slide, "GRIDFORGE", 0.5, 0.9, 7, 1.4, size=72, bold=True, color=COLOR_ACCENT)
    add_text(slide, "Gaming Platform", 0.5, 2.15, 7, 0.7, size=28, color=COLOR_WHITE)
    add_text(slide, "Team  SnackIt  ·  Builders Lab 2026", 0.5, 2.75, 7, 0.5,
             size=15, italic=True, color=COLOR_LIGHT)

    # divider
    div = slide.shapes.add_shape(1, Inches(0.5), Inches(3.35), Inches(5.5), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = COLOR_ACCENT2; div.line.fill.background()

    # Team members cards
    members = [
        ("Danaila Mihai Teodor",  "Frontend Engineer"),
        ("George Bot",             "SPA / SDK Engineer"),
        ("Maria Isabel Oprea",     "QA Engineer"),
        ("Vlad Bontas",            "DevOps Engineer"),
        ("Robert Blehoianu",       "E-commerce / Backend"),
    ]
    for i, (name, role) in enumerate(members):
        x = 0.45
        y = 3.6 + i * 0.67
        add_rect(slide, x, y, 5.3, 0.58, fill=COLOR_CARD, line=COLOR_ACCENT2)
        add_text(slide, name, x+0.15, y+0.06, 3.2, 0.3, size=13, bold=True, color=COLOR_WHITE)
        add_text(slide, role, x+0.15, y+0.3, 3.2, 0.25, size=11, color=COLOR_ACCENT)

    # logo area right
    logo_box = slide.shapes.add_shape(1, Inches(8.8), Inches(2.2), Inches(3.8), Inches(3.5))
    logo_box.fill.solid(); logo_box.fill.fore_color.rgb = COLOR_CARD
    logo_box.line.color.rgb = COLOR_ACCENT; logo_box.line.width = Pt(2)

    add_text(slide, "⬡", 9.4, 2.5, 2.8, 2.0, size=100, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, "G R I D F O R G E", 8.8, 4.8, 3.8, 0.5, size=13, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # slide number
    add_text(slide, "01", 12.8, 7.0, 0.5, 0.4, size=11, color=COLOR_ACCENT2, align=PP_ALIGN.RIGHT)

slide_1()

# ─────────────────────────────────────────────────────────
# SLIDE 2 — What Is GRIDFORGE
# ─────────────────────────────────────────────────────────
def slide_2():
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_accent_bar(slide, t=0.82)

    add_text(slide, "WHAT IS GRIDFORGE?", 0.5, 0.25, 12, 0.7, size=36, bold=True, color=COLOR_ACCENT)
    add_text(slide, "A production-grade gaming platform built in 48 hours — backend, frontend, SDK, tests & infra",
             0.5, 0.9, 12, 0.55, size=16, italic=True, color=COLOR_LIGHT)

    # 3 pillar cards
    pillars = [
        ("⚙  Platform", "Complete REST API powering player accounts, game management, item economy, leaderboards and an ad system — all secured with JWT authentication."),
        ("🖥  Admin CMS", "React + TypeScript admin panel giving non-technical operators full control: create games, manage versions, track players, handle inventory & gifts."),
        ("📦  Developer SDK", "Python client library so any game studio integrates in < 5 minutes. 26 high-level methods. No raw HTTP required."),
    ]
    for i, (title, body) in enumerate(pillars):
        x = 0.4 + i * 4.3
        add_rect(slide, x, 1.6, 4.0, 3.5, fill=COLOR_CARD, line=COLOR_ACCENT)
        add_text(slide, title, x+0.2, 1.75, 3.6, 0.5, size=16, bold=True, color=COLOR_ACCENT)
        div = slide.shapes.add_shape(1, Inches(x+0.2), Inches(2.35), Inches(3.5), Inches(0.03))
        div.fill.solid(); div.fill.fore_color.rgb = COLOR_ACCENT2; div.line.fill.background()
        add_text(slide, body, x+0.2, 2.45, 3.6, 2.5, size=12.5, color=COLOR_LIGHT)

    # tech tags
    tech = ["Laravel PHP", "PostgreSQL", "React + TypeScript", "Python SDK", "JWT Auth", "Docker / K8s"]
    for i, tag in enumerate(tech):
        add_tag(slide, tag, 0.5 + i * 2.15, 5.4)

    add_text(slide, "02", 12.8, 7.0, 0.5, 0.4, size=11, color=COLOR_ACCENT2, align=PP_ALIGN.RIGHT)

slide_2()

# ─────────────────────────────────────────────────────────
# SLIDE 3 — What GRIDFORGE Does
# ─────────────────────────────────────────────────────────
def slide_3():
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_accent_bar(slide, t=0.82)

    add_text(slide, "WHAT GRIDFORGE DOES", 0.5, 0.25, 12, 0.7, size=36, bold=True, color=COLOR_ACCENT)
    add_text(slide, "Six interconnected subsystems delivering a complete gaming experience",
             0.5, 0.9, 12, 0.45, size=15, italic=True, color=COLOR_LIGHT)

    features = [
        ("👤  Player Management",    "Register, authenticate, and manage player profiles.\nRole-based access: Player vs Developer.\nStatus control: active / suspended / banned."),
        ("🎮  Game Lifecycle",        "Create games and versioned releases.\nActivate / deactivate specific versions.\nDeveloper-only write access to protect production."),
        ("🏆  Scoring & Leaderboards","Submit scores (keeps personal best only).\nRanked leaderboard per game.\nPersonal rank endpoint for quick lookups."),
        ("🎁  Item Economy & Gifts",  "Define per-game item catalog with rarities.\nManage player inventories.\nFull gift state machine: pending → accepted/rejected."),
        ("📢  Ad System",             "Ad lifecycle: idle → ad_pending → ad_playing.\nSubscribed players bypass ads automatically.\nPlatform handshake for SDK health checks."),
        ("🔗  SDK & Integration",     "pip install gridforge — zero raw HTTP.\nCopy-paste demo scripts included.\nOpenAPI spec + Swagger UI for documentation."),
    ]

    for i, (title, body) in enumerate(features):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.3
        y = 1.5 + row * 2.7
        add_rect(slide, x, y, 4.0, 2.5, fill=COLOR_CARD, line=COLOR_ACCENT2)
        add_text(slide, title, x+0.15, y+0.1, 3.7, 0.4, size=13, bold=True, color=COLOR_ACCENT)
        add_text(slide, body, x+0.15, y+0.55, 3.7, 1.85, size=11, color=COLOR_LIGHT)

    add_text(slide, "03", 12.8, 7.0, 0.5, 0.4, size=11, color=COLOR_ACCENT2, align=PP_ALIGN.RIGHT)

slide_3()

# ─────────────────────────────────────────────────────────
# PER-PERSON SLIDE HELPER
# ─────────────────────────────────────────────────────────

def person_slide_a(slide_num, name, role, emoji, tagline, highlights, tech_tags):
    """First slide for a person — intro + highlights."""
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)

    # left accent strip
    strip = slide.shapes.add_shape(1, 0, 0, Inches(0.18), prs.slide_height)
    strip.fill.solid(); strip.fill.fore_color.rgb = COLOR_ACCENT; strip.line.fill.background()

    add_accent_bar(slide, t=0.82)

    # avatar circle (simulated)
    avatar = slide.shapes.add_shape(9, Inches(0.5), Inches(0.3), Inches(1.5), Inches(1.5))
    avatar.fill.solid(); avatar.fill.fore_color.rgb = COLOR_ACCENT2
    avatar.line.color.rgb = COLOR_ACCENT; avatar.line.width = Pt(2)
    add_text(slide, emoji, 0.5, 0.3, 1.5, 1.5, size=42, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, name,  2.2, 0.3, 9, 0.65, size=32, bold=True, color=COLOR_WHITE)
    add_text(slide, role,  2.2, 0.9, 9, 0.45, size=18, color=COLOR_ACCENT, bold=True)
    add_text(slide, f'"{tagline}"', 2.2, 1.3, 9, 0.45, size=13, italic=True, color=COLOR_LIGHT)

    # divider
    div = slide.shapes.add_shape(1, Inches(0.35), Inches(1.9), Inches(12.6), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = COLOR_ACCENT2; div.line.fill.background()

    # highlight cards
    for i, (htitle, hbody) in enumerate(highlights):
        col = i % 2
        row = i // 2
        x = 0.35 + col * 6.45
        y = 2.05 + row * 2.35
        add_rect(slide, x, y, 6.1, 2.15, fill=COLOR_CARD, line=COLOR_ACCENT)
        add_text(slide, htitle, x+0.18, y+0.1, 5.7, 0.4, size=13, bold=True, color=COLOR_ACCENT)
        add_text(slide, hbody,  x+0.18, y+0.55, 5.7, 1.5, size=11.5, color=COLOR_LIGHT)

    # tech tags
    for i, tag in enumerate(tech_tags):
        add_tag(slide, tag, 0.35 + i * 2.0, 6.85, color=COLOR_ACCENT2)

    add_text(slide, f"0{slide_num}", 12.8, 7.0, 0.5, 0.4, size=11, color=COLOR_ACCENT2, align=PP_ALIGN.RIGHT)


def person_slide_b(slide_num, name, role_short, emoji, key_points, metrics):
    """Second slide for a person — key contributions + metrics."""
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)

    strip = slide.shapes.add_shape(1, 0, 0, Inches(0.18), prs.slide_height)
    strip.fill.solid(); strip.fill.fore_color.rgb = COLOR_ACCENT2; strip.line.fill.background()

    add_accent_bar(slide, t=0.82, color=COLOR_ACCENT2)

    add_text(slide, f"{emoji}  {name}", 0.35, 0.2, 10, 0.65, size=28, bold=True, color=COLOR_WHITE)
    add_text(slide, role_short, 0.35, 0.78, 10, 0.4, size=15, color=COLOR_ACCENT, bold=True)

    div = slide.shapes.add_shape(1, Inches(0.35), Inches(1.28), Inches(12.6), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = COLOR_ACCENT; div.line.fill.background()

    # left: bullet list
    add_rect(slide, 0.35, 1.45, 7.8, 5.15, fill=COLOR_CARD, line=COLOR_ACCENT)
    add_text(slide, "Key Contributions", 0.55, 1.55, 7.4, 0.4, size=14, bold=True, color=COLOR_ACCENT)
    y = 2.05
    for point in key_points:
        add_text(slide, f"▸  {point}", 0.55, y, 7.4, 0.35, size=12, color=COLOR_LIGHT)
        y += 0.38

    # right: metric boxes
    add_text(slide, "Impact", 8.4, 1.45, 4.6, 0.4, size=14, bold=True, color=COLOR_ACCENT2)
    for i, (label, value) in enumerate(metrics):
        my = 1.95 + i * 1.05
        add_rect(slide, 8.3, my, 4.6, 0.88, fill=COLOR_CARD2, line=COLOR_ACCENT2)
        add_text(slide, value, 8.45, my+0.02, 4.3, 0.48, size=26, bold=True, color=COLOR_ACCENT)
        add_text(slide, label,  8.45, my+0.5,  4.3, 0.35, size=11, color=COLOR_LIGHT)

    add_text(slide, f"0{slide_num}", 12.8, 7.0, 0.5, 0.4, size=11, color=COLOR_ACCENT2, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────
# SLIDES 4–5  Danaila Mihai Teodor — Frontend
# ─────────────────────────────────────────────────────────
person_slide_a(
    4,
    "Danaila Mihai Teodor",
    "Frontend Engineer",
    "🖥",
    "Turned raw API endpoints into a polished admin experience",
    [
        ("React + TypeScript CMS",
         "Built the full admin panel from scratch using React, TypeScript and Vite.\nAuthentication screens, dashboard navigation, and role-based views for both players and developers."),
        ("Game & Version Management",
         "Create / edit games and versioned releases directly from the UI.\nGameCard and VersionForm components with live API integration."),
        ("Player & Inventory Screens",
         "Browse player profiles, view progress per game, and manage item inventories — all from a single responsive dashboard."),
        ("TypeScript Data Models",
         "Defined typed interfaces for every domain entity (Player, Game, Gift, Leaderboard) ensuring compile-time safety across the whole frontend."),
    ],
    ["React", "TypeScript", "Vite", "Context API", "REST Integration"]
)

person_slide_b(
    5,
    "Danaila Mihai Teodor",
    "Frontend Engineer — CMS & Admin Panel",
    "🖥",
    [
        "Implemented login / signup with context-based auth state management",
        "Dashboard with multi-section navigation: Games, Players, Inventory, Gifts",
        "GameCard, ItemForm, VersionForm — reusable typed UI components",
        "Role-aware views: different experience for Player vs Developer",
        "Live data fetching from 40+ API endpoints with error handling",
        "Gift management UI with pending / accepted / rejected state display",
        "Responsive layouts across all major screen breakpoints",
        "18 commits — most frontend contributions on the team",
    ],
    [
        ("Commits (frontend owner)", "18"),
        ("API endpoints integrated", "40 +"),
        ("TypeScript models defined", "10 +"),
        ("Dashboard sections", "4"),
        ("Component library size", "15 +"),
    ]
)

# ─────────────────────────────────────────────────────────
# SLIDES 6–7  George Bot — SPA / SDK
# ─────────────────────────────────────────────────────────
person_slide_a(
    6,
    "George Bot",
    "SPA Software Engineer — SDK & Integration",
    "🐍",
    "First published score in under 5 minutes — zero API docs required",
    [
        ("Python SDK — gridforge",
         "Developed the full Python client library with 26 high-level methods covering auth, profiles, scoring, leaderboards, inventory, gifts and ads."),
        ("Typed Exception Hierarchy",
         "AuthError, ForbiddenError, NotFoundError, ConflictError, ValidationError, APIError — every failure mode is catchable and documented."),
        ("Demo & Game Scripts",
         "demo.py: full feature showcase.\ngame.py: interactive number-guessing game that calls the real platform — a living integration example."),
        ("Cookie + Bearer Token Auth",
         "Dual auth storage (cookie & header) so the SDK works transparently in browser-based tools and CLI scripts alike."),
    ],
    ["Python", "requests", "JWT", "OpenAPI", "pip package"]
)

person_slide_b(
    7,
    "George Bot",
    "SPA Software Engineer — SDK & Game Integration",
    "🐍",
    [
        "26 SDK methods: register, login, submit_score, get_leaderboard, send_gift, get_ad_status …",
        "Typed exceptions so errors never surface as raw HTTP status codes",
        "demo.py walks through the full auth → score → leaderboard → gift flow",
        "game.py — interactive guessing game running against the live platform API",
        "Cookie handling added to support browser-based and headless clients",
        "5-step integration path: pip install → first score in < 5 minutes",
        "API docs auto-generation wired into the project pipeline",
        "18 total commits across SDK and integration work",
    ],
    [
        ("SDK methods shipped", "26"),
        ("Exception types defined", "6"),
        ("Lines of client code", "400 +"),
        ("Integration steps to first score", "5"),
        ("Total commits", "18"),
    ]
)

# ─────────────────────────────────────────────────────────
# SLIDES 8–9  Maria Isabel Oprea — QA
# ─────────────────────────────────────────────────────────
person_slide_a(
    8,
    "Maria Isabel Oprea",
    "QA Automation Engineer",
    "🧪",
    "Mock first, real API second — tests that evolve with the platform",
    [
        ("Automation Strategy",
         "Defined a written QA strategy prioritising 3 critical flows: authentication, gift lifecycle, and leaderboard ranking — with justification for each choice."),
        ("Phase 1 — Mock Tests",
         "Complete pytest suite using the responses library.\nSecurity: unauthenticated blocking, invalid token rejection.\nFunctional: score submission, profile update, gift state transitions."),
        ("Phase 2 — Real API Tests",
         "Same test code, decorator removed — mock layer replaced with live API calls without rewriting tests. Clean migration path validated."),
        ("Negative & E2E Testing",
         "Full user flows: registration → auth → action → validation.\nNegative cases: duplicate registration, input validation errors, 401/422 assertions."),
    ],
    ["pytest", "responses (mock)", "requests", "E2E testing", "API testing"]
)

person_slide_b(
    9,
    "Maria Isabel Oprea",
    "QA Engineer — Test Automation & Strategy",
    "🧪",
    [
        "Written AUTOMATION_STRATEGY.md — 3 critical flows, rationale documented",
        "Security tests: unauthenticated requests return 401, invalid tokens blocked",
        "Validation tests: duplicate registration → 422, malformed input rejected",
        "Gift lifecycle: pending → accepted / rejected state machine verified",
        "Leaderboard ranking assertion: sorted by score, not insertion order",
        "Profile read/write consistency checks across create → update → fetch",
        "Phase 1 (mock) → Phase 2 (real) migration with zero test rewrites",
        "14 commits — full ownership of QA layer",
    ],
    [
        ("Commits (QA owner)", "14"),
        ("Critical flows prioritised", "3"),
        ("Test phases shipped", "2"),
        ("HTTP status codes asserted", "401, 422"),
        ("Test types covered", "Security, E2E, Negative"),
    ]
)

# ─────────────────────────────────────────────────────────
# SLIDES 10–11  Vlad Bontas — DevOps
# ─────────────────────────────────────────────────────────
person_slide_a(
    10,
    "Vlad Bontas",
    "DevOps Engineer",
    "⚙",
    "Reproducible infra from code — zero manual clicks in production",
    [
        ("Kubernetes with kind",
         "Local cluster in Docker (< 1 min spin-up, zero cloud cost) running the full platform stack with Nginx Ingress as the single entry point."),
        ("Terraform + AWS",
         "AWS RDS (MySQL) provisioned with Terraform for persistent state — compute is disposable, data is not. Infrastructure as code, fully version-controlled."),
        ("GitHub Actions CI/CD",
         "Two pipelines: unit tests → Docker build → Trivy security scan → Helm deploy.\nSelf-hosted runner on EC2 for full pipeline control."),
        ("Observability & Load Testing",
         "Prometheus + Grafana for live metrics.\nk6 load tests validating HPA auto-scaling from 1 → 10 replicas under sustained load."),
    ],
    ["Kubernetes", "Terraform", "Helm", "Docker", "GitHub Actions", "Prometheus", "Grafana", "k6"]
)

person_slide_b(
    11,
    "Vlad Bontas",
    "DevOps Engineer — Infrastructure & CI/CD",
    "⚙",
    [
        "kind (Kubernetes in Docker) cluster — full K8s API, zero EKS overhead",
        "AWS RDS separated from compute: disposable pods, persistent data",
        "Nginx Ingress: single auditable entry point, ClusterIP for internal routing",
        "Helm charts for all application deployments (devops/helm/vector-api)",
        "Trivy security scan in every CI pipeline — containers checked before deploy",
        "HPA auto-scaling validated: 1 → 10 replicas under k6 sustained load",
        "Fixed 3 real production issues: pymysql dialect, python3-venv, HPA metrics",
        "DECISIONS-DEVOPS.md: documented rationale + troubleshooting log",
    ],
    [
        ("CI/CD pipelines", "2"),
        ("Real issues debugged & fixed", "3"),
        ("HPA max replicas scaled to", "10"),
        ("IaC tools used", "Terraform + Helm"),
        ("Observability stack", "Prometheus + Grafana"),
    ]
)

# ─────────────────────────────────────────────────────────
# SLIDES 12–13  Robert Blehoianu — Backend / E-commerce
# ─────────────────────────────────────────────────────────
person_slide_a(
    12,
    "Robert Blehoianu",
    "E-commerce & Backend Engineer",
    "🔧",
    "API-first: spec written before the first line of code",
    [
        ("Laravel REST API",
         "Architected the full backend with domain-driven structure: User, Game and Transaction pillars — each with its own Models, Controllers and migrations."),
        ("OpenAPI Spec (40 + endpoints)",
         "Defined the complete API contract in openapi.yaml before implementation.\nBearer JWT security scheme, all request/response schemas specified upfront."),
        ("Item Economy & Gift System",
         "Per-game item catalog with rarities, player inventories with unique constraints, and a full gift state machine (pending → accepted / rejected)."),
        ("UUID-first Database Design",
         "All primary keys are UUIDs for security and portability.\nCASCADE deletes, unique constraints, and Neon PostgreSQL with connection pooling."),
    ],
    ["Laravel PHP 8.3", "PostgreSQL", "JWT", "OpenAPI", "UUID", "Swagger UI"]
)

person_slide_b(
    13,
    "Robert Blehoianu",
    "E-commerce & Backend Engineer — REST API Architecture",
    "🔧",
    [
        "Designed 3-pillar domain structure: User / Game / Transaction controllers",
        "JWT stateless auth: register, login, logout, cookie fallback for browsers",
        "Role-based access: Developer-only write endpoints on games, items, versions",
        "Score upsert: only the personal best is kept per player per game",
        "Gift state machine: pending → accepted / rejected (irreversible transitions)",
        "Ad system state machine: idle → ad_pending → ad_playing + subscription bypass",
        "10 database migrations, all with UUID PKs and proper CASCADE constraints",
        "21 commits — highest commit count on the team",
    ],
    [
        ("Commits (most on team)", "21"),
        ("API endpoints designed", "40 +"),
        ("Database migrations", "10"),
        ("Domain controllers", "9"),
        ("Decisions documented", "9"),
    ]
)

# ─────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────
out = r"C:\Users\Ciocanc246\Hackaton\team-7\SnackIt_GRIDFORGE.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
