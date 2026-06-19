from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x09, 0x09, 0x18)
CARD     = RGBColor(0x12, 0x12, 0x28)
HDR_HUB  = RGBColor(0x5B, 0x21, 0xB6)   # purple  – core entities
HDR_CHLD = RGBColor(0x0C, 0x2A, 0x55)   # navy    – child / junction tables
CYAN     = RGBColor(0x00, 0xE5, 0xFF)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x8A, 0xA4, 0xC8)
DIM      = RGBColor(0x55, 0x60, 0x7A)
GOLD     = RGBColor(0xFF, 0xCC, 0x00)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
LINE_C   = RGBColor(0x2A, 0x3A, 0x6A)   # subtle connector color

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Primitive helpers ────────────────────────────────────────────────────────
def box(l, t, w, h, fill=CARD, border=None, bw=1.2):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else:       s.line.fill.background()
    return s

def label(l, t, w, h, text, size=10, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def line(x1, y1, x2, y2, color=LINE_C, w=1.0):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)

# ── Background ───────────────────────────────────────────────────────────────
box(0, 0, 13.33, 7.5, fill=BG)

# subtle grid lines
for gx in [2.65, 5.3, 7.95, 10.6]:
    b = slide.shapes.add_shape(1, Inches(gx), 0, Pt(1), prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x2E)
    b.line.fill.background()

# ── Title bar ────────────────────────────────────────────────────────────────
box(0, 0, 13.33, 0.54, fill=RGBColor(0x06, 0x06, 0x14))
label(0.28, 0.06, 9.5, 0.42, "GRIDFORGE  ·  Database Entity Relationship Diagram",
      size=19, bold=True, color=CYAN)
label(9.8, 0.12, 3.3, 0.32, "9 tables  ·  UUID PKs  ·  PostgreSQL (Neon)",
      size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# cyan underline + purple accent
acc1 = slide.shapes.add_shape(1, 0, Inches(0.54), prs.slide_width, Inches(0.030))
acc1.fill.solid(); acc1.fill.fore_color.rgb = CYAN; acc1.line.fill.background()
acc2 = slide.shapes.add_shape(1, 0, Inches(0.570), prs.slide_width, Inches(0.016))
acc2.fill.solid(); acc2.fill.fore_color.rgb = PURPLE; acc2.line.fill.background()

# ── Entity builder ───────────────────────────────────────────────────────────
HDR  = 0.285   # header height
ROW  = 0.200   # field row height
PAD  = 0.055   # bottom padding inside box

FIELD_COLORS = {
    'PK': (GOLD,  "⬡"),
    'FK': (CYAN,  "⟶"),
    '':   (MUTED, "  "),
}

def entity(name, fields, x, y, w=2.20, hdr=HDR_CHLD):
    h = HDR + len(fields) * ROW + PAD
    # shadow
    box(x + 0.04, y + 0.04, w, h, fill=RGBColor(0x05, 0x05, 0x10))
    # main card
    box(x, y, w, h, fill=CARD, border=CYAN, bw=1.0)
    # header fill
    box(x, y, w, HDR, fill=hdr)
    # header divider
    d = slide.shapes.add_shape(1, Inches(x), Inches(y + HDR), Inches(w), Inches(0.016))
    d.fill.solid(); d.fill.fore_color.rgb = CYAN; d.line.fill.background()
    # header text
    label(x + 0.1, y + 0.048, w - 0.2, HDR - 0.06, name,
          size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # fields
    for i, (fname, ftype, tag) in enumerate(fields):
        fy    = y + HDR + 0.018 + i * ROW
        color, sym = FIELD_COLORS[tag]
        label(x + 0.10, fy, w * 0.66, ROW, f"{sym} {fname}", size=8.5, color=color)
        label(x + w * 0.66, fy, w * 0.31, ROW, ftype,
              size=8, color=DIM, align=PP_ALIGN.RIGHT)
    return {'l': x, 't': y, 'r': x + w, 'b': y + h,
            'mx': x + w / 2, 'my': y + h / 2, 'w': w, 'h': h}

# ── Layout: two rows ─────────────────────────────────────────────────────────
#   Row 1  y = 0.64   roles | players | player_profiles | games | game_versions
#   Row 2  y = 3.30   gifts | player_progress | player_items | items

E = {}

# Row 1
E['roles']           = entity('roles', [
    ('id',   'UUID',   'PK'),
    ('name', 'string', ''),
], x=0.22, y=0.64, w=2.05, hdr=HDR_HUB)

E['players']         = entity('players', [
    ('id',            'UUID',      'PK'),
    ('username',      'string',    ''),
    ('email',         'string',    ''),
    ('password_hash', 'string',    ''),
    ('role_id',       'UUID',      'FK'),
    ('status',        'enum',      ''),
    ('created_at',    'timestamp', ''),
], x=2.55, y=0.64, w=2.35, hdr=HDR_HUB)

E['player_profiles'] = entity('player_profiles', [
    ('id',            'UUID',      'PK'),
    ('player_id',     'UUID',      'FK'),
    ('display_name',  'string',    ''),
    ('subscribed',    'boolean',   ''),
    ('ad_state',      'enum',      ''),
    ('last_ad_shown', 'timestamp', ''),
], x=5.18, y=0.64, w=2.40, hdr=HDR_CHLD)

E['games']           = entity('games', [
    ('id',          'UUID',      'PK'),
    ('name',        'string',    ''),
    ('description', 'text',      ''),
    ('created_at',  'timestamp', ''),
], x=7.90, y=0.64, w=2.20, hdr=HDR_HUB)

E['game_versions']   = entity('game_versions', [
    ('id',          'UUID',      'PK'),
    ('game_id',     'UUID',      'FK'),
    ('version',     'string',    ''),
    ('is_active',   'boolean',   ''),
    ('released_at', 'timestamp', ''),
], x=10.42, y=0.64, w=2.65, hdr=HDR_CHLD)

# Row 2
E['gifts']           = entity('gifts', [
    ('id',           'UUID',      'PK'),
    ('sender_id',    'UUID',      'FK'),
    ('recipient_id', 'UUID',      'FK'),
    ('item_id',      'UUID',      'FK'),
    ('status',       'enum',      ''),
    ('sent_at',      'timestamp', ''),
], x=0.22, y=3.30, w=2.40, hdr=HDR_CHLD)

E['player_progress'] = entity('player_progress', [
    ('id',         'UUID',      'PK'),
    ('player_id',  'UUID',      'FK'),
    ('game_id',    'UUID',      'FK'),
    ('score',      'integer',   ''),
    ('updated_at', 'timestamp', ''),
], x=2.95, y=3.30, w=2.35, hdr=HDR_CHLD)

E['player_items']    = entity('player_items', [
    ('id',        'UUID',    'PK'),
    ('player_id', 'UUID',    'FK'),
    ('item_id',   'UUID',    'FK'),
    ('quantity',  'integer', ''),
], x=5.60, y=3.30, w=2.20, hdr=HDR_CHLD)

E['items']           = entity('items', [
    ('id',          'UUID',   'PK'),
    ('game_id',     'UUID',   'FK'),
    ('name',        'string', ''),
    ('description', 'text',   ''),
    ('rarity',      'enum',   ''),
], x=8.12, y=3.30, w=2.20, hdr=HDR_CHLD)

# ── Connectors ───────────────────────────────────────────────────────────────
# Draw connections BEFORE entities so they appear behind — but since we
# already drew entities, we draw connectors last (they sit on top which is fine).

REL = CYAN      # relationship line
REL2 = GREEN    # junction table lines

# Row 1 horizontal
# roles ──▶ players  (role_id)
line(E['roles']['r'],  E['roles']['my'],
     E['players']['l'], E['players']['t'] + HDR + ROW * 3.5,
     color=REL, w=1.4)

# players ──▶ player_profiles  (player_id)
line(E['players']['r'], E['players']['t'] + HDR + ROW * 0.5,
     E['player_profiles']['l'], E['player_profiles']['t'] + HDR + ROW * 0.5,
     color=REL, w=1.4)

# games ──▶ game_versions  (game_id)
line(E['games']['r'], E['games']['my'],
     E['game_versions']['l'], E['game_versions']['my'],
     color=REL, w=1.4)

# Row 1 ──▶ Row 2 (vertical / diagonal)
# players ──▶ gifts  (sender_id)
line(E['players']['l'] + 0.3, E['players']['b'],
     E['gifts']['r'] - 0.3,   E['gifts']['t'],
     color=PURPLE, w=1.2)

# players ──▶ player_progress  (player_id)
line(E['players']['mx'], E['players']['b'],
     E['player_progress']['l'] + 0.5, E['player_progress']['t'],
     color=REL, w=1.4)

# players ──▶ player_items  (player_id)
line(E['players']['r'] - 0.3, E['players']['b'],
     E['player_items']['l'] + 0.4, E['player_items']['t'],
     color=REL2, w=1.2)

# games ──▶ player_progress  (game_id)
line(E['games']['l'] + 0.5, E['games']['b'],
     E['player_progress']['r'] - 0.4, E['player_progress']['t'],
     color=REL, w=1.4)

# games ──▶ items  (game_id)
line(E['games']['mx'], E['games']['b'],
     E['items']['mx'], E['items']['t'],
     color=REL, w=1.4)

# Row 2 horizontal
# items ──▶ player_items  (item_id)
line(E['items']['l'], E['items']['my'],
     E['player_items']['r'], E['player_items']['my'],
     color=REL2, w=1.2)

# items ──▶ gifts  (item_id)
line(E['gifts']['r'] - 0.2, E['gifts']['b'] - 0.3,
     E['items']['l'], E['items']['b'] - 0.3,
     color=PURPLE, w=1.1)

# ── Cardinality badges  (tiny "1" and "N" labels near line ends) ─────────────
def badge(x, y, txt_val, color=MUTED):
    label(x, y, 0.22, 0.20, txt_val, size=8, bold=True, color=color)

# roles → players
badge(E['roles']['r'] + 0.02,    E['roles']['my'] - 0.22,    "1",  CYAN)
badge(E['players']['l'] - 0.18,  E['players']['t'] + HDR + ROW * 3.1, "N", GOLD)

# players → player_profiles
badge(E['players']['r'] + 0.02,  E['players']['t'] + HDR - 0.05, "1", CYAN)
badge(E['player_profiles']['l'] - 0.18, E['player_profiles']['t'] + HDR + 0.0, "1", GOLD)

# games → game_versions
badge(E['games']['r'] + 0.02,         E['games']['my'] - 0.20,      "1", CYAN)
badge(E['game_versions']['l'] - 0.18, E['game_versions']['my'] - 0.20, "N", GOLD)

# players → player_progress
badge(E['players']['mx'] + 0.02,      E['players']['b'] + 0.01,           "1", CYAN)
badge(E['player_progress']['l'] + 0.3, E['player_progress']['t'] - 0.22, "N", GOLD)

# games → items
badge(E['games']['mx'] + 0.02, E['games']['b'] + 0.01,  "1", CYAN)
badge(E['items']['mx'] + 0.02, E['items']['t'] - 0.22,  "N", GOLD)

# items → player_items
badge(E['items']['l'] - 0.18,       E['items']['my'] - 0.20,        "1", CYAN)
badge(E['player_items']['r'] + 0.02, E['player_items']['my'] - 0.20, "N", GOLD)

# ── Legend ───────────────────────────────────────────────────────────────────
LY = 6.97
box(0.22, LY, 12.85, 0.38, fill=RGBColor(0x0C, 0x0C, 0x22),
    border=RGBColor(0x20, 0x20, 0x44), bw=0.8)

# legend items
items_leg = [
    (GOLD,   "⬡  Primary Key (UUID)"),
    (CYAN,   "⟶  Foreign Key"),
    (CYAN,   "─── 1-to-Many"),
    (GREEN,  "─── Junction table"),
    (PURPLE, "─── Player-gift link"),
    (MUTED,  "1 / N  Cardinality"),
]
for i, (col, txt_val) in enumerate(items_leg):
    label(0.4 + i * 2.1, LY + 0.09, 2.0, 0.24, txt_val, size=8.5, color=col)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\Ciocanc246\Hackaton\team-7\GRIDFORGE_ERD.pptx"
prs.save(out)
print(f"Saved: {out}")
